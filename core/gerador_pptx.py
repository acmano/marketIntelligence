"""
Gerador do relatorio PowerPoint semanal do Market Intelligence.
Usa python-pptx com identidade visual Lorenzetti.
Compativel com Python 3.8+
"""

import os
from datetime import datetime, timedelta
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# Cores Lorenzetti
AZUL_ESCURO  = RGBColor(0x1B, 0x3A, 0x6B)
AZUL_MEDIO   = RGBColor(0x2E, 0x6D, 0xB4)
AZUL_CLARO   = RGBColor(0xD6, 0xE4, 0xF0)
BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
CINZA        = RGBColor(0x55, 0x55, 0x55)
PRETO        = RGBColor(0x22, 0x22, 0x22)

# Dimensoes slide 16:9
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _rgb(r, g, b):
    return RGBColor(r, g, b)


def _fill_solid(shape, cor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor


def _add_textbox(slide, left, top, width, height,
                 texto, tamanho=18, cor=PRETO, bold=False,
                 alinhamento=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = alinhamento
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(tamanho)
    run.font.color.rgb = cor
    run.font.bold = bold
    return txb


def _add_rect(slide, left, top, width, height, cor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    _fill_solid(shape, cor)
    shape.line.fill.background()
    return shape


def _slide_capa(prs, semana_ref: str, data_geracao: str):
    layout = prs.slide_layouts[6]  # em branco
    slide = prs.slides.add_slide(layout)

    # Fundo azul escuro
    _add_rect(slide, 0, 0, 13.33, 7.5, AZUL_ESCURO)

    # Faixa inferior
    _add_rect(slide, 0, 6.0, 13.33, 1.5, AZUL_MEDIO)

    # Titulo
    _add_textbox(slide, 0.8, 1.5, 11.0, 1.5,
                 "MARKET INTELLIGENCE", 44, BRANCO, True, PP_ALIGN.CENTER)

    # Subtitulo
    _add_textbox(slide, 0.8, 3.0, 11.0, 0.8,
                 "Assessoria a Alta Direcao | Lorenzetti S.A.", 22, AZUL_CLARO,
                 False, PP_ALIGN.CENTER)

    # Semana
    _add_textbox(slide, 0.8, 3.9, 11.0, 0.6,
                 semana_ref, 20, BRANCO, False, PP_ALIGN.CENTER)

    # Data de geracao
    _add_textbox(slide, 0.8, 6.1, 11.0, 0.5,
                 f"Gerado automaticamente em {data_geracao} | Uso interno confidencial",
                 14, BRANCO, False, PP_ALIGN.CENTER)


def _slide_titulo_secao(prs, titulo: str, subtitulo: str = ""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    _add_rect(slide, 0, 0, 13.33, 7.5, AZUL_CLARO)
    _add_rect(slide, 0, 0, 0.2, 7.5, AZUL_ESCURO)

    _add_textbox(slide, 0.6, 2.5, 12.0, 1.5,
                 titulo, 36, AZUL_ESCURO, True, PP_ALIGN.LEFT)

    if subtitulo:
        _add_textbox(slide, 0.6, 4.2, 12.0, 0.8,
                     subtitulo, 20, CINZA, False, PP_ALIGN.LEFT)


def _slide_conteudo(prs, titulo_slide: str, corpo: str, categoria_cor=None):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Cabecalho
    cor_cab = categoria_cor or AZUL_ESCURO
    _add_rect(slide, 0, 0, 13.33, 1.1, cor_cab)
    _add_textbox(slide, 0.3, 0.15, 12.5, 0.8,
                 titulo_slide, 24, BRANCO, True, PP_ALIGN.LEFT)

    # Linha separadora
    _add_rect(slide, 0, 1.1, 13.33, 0.05, AZUL_MEDIO)

    # Corpo do texto
    _add_textbox(slide, 0.4, 1.3, 12.5, 5.8,
                 corpo, 16, PRETO, False, PP_ALIGN.LEFT)

    # Rodape
    _add_rect(slide, 0, 7.1, 13.33, 0.4, AZUL_CLARO)
    _add_textbox(slide, 0.3, 7.15, 12.5, 0.3,
                 "Lorenzetti S.A. | Market Intelligence | Uso interno confidencial",
                 10, CINZA, False, PP_ALIGN.CENTER)


def _slide_artigos(prs, titulo_slide: str, artigos: List[Dict[str, Any]], cor=None):
    """Slide com lista de artigos e seus resumos."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    cor_cab = cor or AZUL_ESCURO
    _add_rect(slide, 0, 0, 13.33, 1.1, cor_cab)
    _add_textbox(slide, 0.3, 0.15, 12.5, 0.8,
                 titulo_slide, 24, BRANCO, True, PP_ALIGN.LEFT)
    _add_rect(slide, 0, 1.1, 13.33, 0.05, AZUL_MEDIO)

    top = 1.25
    for i, a in enumerate(artigos[:3], 1):
        data = a.get("data_publicacao")
        data_str = data.strftime("%d/%m/%Y") if data else ""
        fonte = a.get("fonte_nome", "")
        titulo = a.get("titulo", "")[:100]
        resumo = (a.get("resumo_pt") or "")[:200]

        _add_textbox(slide, 0.4, top, 12.5, 0.35,
                     f"{i}. [{fonte}] {titulo}",
                     13, AZUL_ESCURO, True, PP_ALIGN.LEFT)
        top += 0.35

        if resumo:
            _add_textbox(slide, 0.7, top, 12.0, 0.55,
                         resumo, 12, CINZA, False, PP_ALIGN.LEFT)
            top += 0.58

        _add_rect(slide, 0.4, top, 12.5, 0.02, AZUL_CLARO)
        top += 0.1

    _add_rect(slide, 0, 7.1, 13.33, 0.4, AZUL_CLARO)
    _add_textbox(slide, 0.3, 7.15, 12.5, 0.3,
                 "Lorenzetti S.A. | Market Intelligence | Uso interno confidencial",
                 10, CINZA, False, PP_ALIGN.CENTER)


def _slide_fontes(prs, fontes: List[str], periodo: str):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    _add_rect(slide, 0, 0, 13.33, 1.1, AZUL_ESCURO)
    _add_textbox(slide, 0.3, 0.15, 12.5, 0.8,
                 "Fontes Consultadas", 24, BRANCO, True, PP_ALIGN.LEFT)

    _add_textbox(slide, 0.4, 1.2, 12.5, 0.5,
                 f"Periodo: {periodo}", 16, CINZA, False, PP_ALIGN.LEFT)

    corpo = "\n".join(f"  - {f}" for f in fontes)
    _add_textbox(slide, 0.4, 1.8, 12.5, 5.0,
                 corpo, 16, PRETO, False, PP_ALIGN.LEFT)

    _add_rect(slide, 0, 7.1, 13.33, 0.4, AZUL_CLARO)
    _add_textbox(slide, 0.3, 7.15, 12.5, 0.3,
                 "Lorenzetti S.A. | Market Intelligence | Uso interno confidencial",
                 10, CINZA, False, PP_ALIGN.CENTER)


def gerar_powerpoint(
    semana_inicio: datetime,
    resumo_executivo: str,
    conteudo_materias: str,
    artigos_materias: List[Dict],
    conteudo_exportacao: str,
    artigos_exportacao: List[Dict],
    conteudo_geopolitica: str,
    artigos_geopolitica: List[Dict],
    alertas: str,
    analise_preditiva: str,
    fontes_ativas: List[str],
    caminho_saida: Optional[str] = None,
) -> str:
    """
    Gera o PowerPoint semanal completo.
    Retorna o caminho do arquivo gerado.
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    semana_fim = semana_inicio + timedelta(days=6)
    semana_ref = f"Semana de {semana_inicio.strftime('%d/%m')} a {semana_fim.strftime('%d/%m/%Y')}"
    data_geracao = datetime.now().strftime("%d/%m/%Y as %H:%M")

    # 1. Capa
    _slide_capa(prs, semana_ref, data_geracao)

    # 2. Sumario executivo
    _slide_conteudo(prs, "Sumario Executivo da Semana", resumo_executivo, AZUL_MEDIO)

    # 3. Materias-primas — analise
    _slide_titulo_secao(prs, "Materias-Primas Plasticas", "Polipropileno | ABS | Nylon | Derivados")
    _slide_artigos(prs, "Principais Noticias — Materias-Primas", artigos_materias, AZUL_ESCURO)
    _slide_conteudo(prs, "Analise — Materias-Primas", conteudo_materias, AZUL_ESCURO)

    # 4. Mercados de exportacao — analise
    _slide_titulo_secao(prs, "Mercados de Exportacao", "Demanda | Regulacao | Concorrencia")
    _slide_artigos(prs, "Principais Noticias — Exportacao", artigos_exportacao, _rgb(0x1A, 0x6B, 0x3A))
    _slide_conteudo(prs, "Analise — Mercados de Exportacao", conteudo_exportacao, _rgb(0x1A, 0x6B, 0x3A))

    # 5. Geopolitica / Macroeconomia
    _slide_titulo_secao(prs, "Geopolitica e Macroeconomia", "Contexto global relevante para a Lorenzetti")
    _slide_artigos(prs, "Destaques — Geopolitica e Economia", artigos_geopolitica, _rgb(0x7B, 0x3A, 0x1B))
    _slide_conteudo(prs, "Analise — Geopolitica", conteudo_geopolitica, _rgb(0x7B, 0x3A, 0x1B))

    # 6. Alertas e riscos
    _slide_conteudo(prs, "Alertas e Riscos Identificados", alertas, _rgb(0xB4, 0x2E, 0x2E))

    # 7. Analise preditiva
    _slide_conteudo(prs, "Analise Preditiva — Proximas Semanas", analise_preditiva, _rgb(0x2E, 0x4D, 0xB4))

    # 8. Fontes
    _slide_fontes(prs, fontes_ativas, semana_ref)

    # Salva
    if not caminho_saida:
        data_str = semana_inicio.strftime("%Y-%m-%d")
        caminho_saida = os.path.join(
            os.environ.get("MI_PROJECT_ROOT", "/home/mano/projetos/datasul/marketIntelligence"),
            "output",
            f"mi_relatorio_{data_str}.pptx"
        )

    os.makedirs(os.path.dirname(caminho_saida), exist_ok=True)
    prs.save(caminho_saida)
    return caminho_saida
